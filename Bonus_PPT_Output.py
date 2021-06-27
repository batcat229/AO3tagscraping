from pptx import Presentation
import xlrd
import codecs

#a ppt auto output code
#originally use for character top 100

prs = Presentation('Char_100.pptx')

workbook = xlrd.open_workbook('Result.xlsx')
worksheet = workbook.sheet_by_name("True")
file=codecs.open("top_100_char.txt","w",'utf-8')

row=0

for row in range(101,0,-1):
    #Extrat items
    rank=str(int(worksheet.cell(row,0).value))
    char_name=str(worksheet.cell(row,1).value)
    hit=str(int(worksheet.cell(row,3).value))
    fandom=str(worksheet.cell(row,4).value)
    fandom_type=str(worksheet.cell(row,5).value)
    sex=str(worksheet.cell(row,6).value)
    color=str(worksheet.cell(row,7).value)
    #Add Slide
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)
    #Write the Slide
    shapes=slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = rank + '.' + char_name
    tf = body_shape.text_frame
    tf.text='Hit: '+hit
    p = tf.add_paragraph()
    p.text = 'Fandom: '+fandom
    p = tf.add_paragraph()
    p.text = 'Fandom Type: '+fandom_type
    p = tf.add_paragraph()
    p.text = 'Sex: '+sex
    p = tf.add_paragraph()
    p.text = 'Color: '+color

prs.save('Char_100.pptx')